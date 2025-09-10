# ============================================================
# EDA Chancado — FUSIÓN (Streamlit Cloud, versión simple y útil)
# Controles: Periodo, Métrica foco, Brecha Δ/%, Outliers, Sensibilidad,
#            ANOVA α, Exportar PPT con notas del presentador
# ============================================================

import io, re
import numpy as np
import pandas as pd
import streamlit as st
import seaborn as sns
import matplotlib.pyplot as plt
from matplotlib.ticker import StrMethodFormatter
from matplotlib.dates import DateFormatter, MonthLocator
from matplotlib import colors as mcolors

import statsmodels.api as sm
from statsmodels.formula.api import ols
from statsmodels.stats.multicomp import pairwise_tukeyhsd, MultiComparison

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

st.set_page_config(page_title="EDA Chancado — FUSIÓN", layout="wide")
sns.set_theme(style="whitegrid")

# ---------- Paleta ----------
PALETTE = ["#328BA1", "#0B5563", "#00AFAA", "#66C7C7", "#BFD8D2", "#003B5C", "#7FB7BE"]
custom_palette = PALETTE[:5]

# ---------- Helpers ----------
def _normalize_colname(c: str) -> str:
    return (c.strip().lower()
            .replace("/", "_")
            .replace(" ", "_")
            .replace("__","_"))

EXPECTED = {
    "Fecha": ["fecha","date"],
    "mineral_procesado_real_t": ["mineral_procesado_real_t","mineral_real_t","ton_real","tons_real"],
    "rendimiento_real_tph": ["rendimiento_real_tph","tph_real","real_tph"],
    "tiempo_operativo_real_h/dia": ["tiempo_operativo_real_h_dia","tiempo_operativo_real_h/dia","horas_reales","h_real"],
    "mineral_procesado_plan_t": ["mineral_procesado_plan_t","mineral_plan_t","ton_plan","tons_plan"],
    "rendimiento_plan_tph": ["rendimiento_plan_tph","tph_plan","plan_tph"],
    "tiempo_operativo_plan_h/dia": ["tiempo_operativo_plan_h_dia","tiempo_operativo_plan_h/dia","horas_plan","h_plan"]
}

def _map_columns(df_cols):
    norm_map = {_normalize_colname(c): c for c in df_cols}
    mapping = {}
    for std, aliases in EXPECTED.items():
        found = None
        for a in aliases + [_normalize_colname(std)]:
            if a in norm_map:
                found = norm_map[a]; break
        mapping[std] = found
    missing = [k for k,v in mapping.items() if v is None]
    if missing:
        raise ValueError(f"Faltan columnas requeridas en el CSV: {missing}")
    return mapping

def _to_num(series: pd.Series) -> pd.Series:
    return pd.to_numeric(
        series.astype(str)
              .str.replace("\u00a0","", regex=False)
              .str.replace(".", "", regex=False)
              .str.replace(",", "", regex=False)
              .str.strip(),
        errors="coerce"
    )

def _parse_fecha_mmddyyyy_with_report(s: pd.Series):
    raw = s.astype(str).str.strip()
    pat_exact = re.compile(r"^\d{2}-\d{2}\.\d{4}$")
    exact_count = int(raw.apply(lambda x: bool(pat_exact.match(x))).sum())

    s_norm = (raw
              .str.replace("\u00a0","", regex=False)
              .str.replace("/", "-", regex=False)
              .str.replace(".", "-", regex=False))

    fechas_strict = pd.to_datetime(s_norm, format="%m-%d-%Y", errors="coerce")
    strict_ok = fechas_strict.notna()
    fechas = fechas_strict.copy()

    if (~strict_ok).any():
        fechas_fallback = pd.to_datetime(s_norm[~strict_ok], errors="coerce", infer_datetime_format=True)
        fechas.loc[~strict_ok] = fechas_fallback

    report = {
        "total": len(s),
        "exact_mm-dd.yyyy": exact_count,
        "parsed": int(fechas.notna().sum()),
        "invalid": int(fechas.isna().sum())
    }
    st.caption(f"🗓️ Validación fechas: {report}")
    return fechas

def load_dataset_v2(df_in: pd.DataFrame) -> pd.DataFrame:
    mapping = _map_columns(df_in.columns)
    out = pd.DataFrame()
    out["Fecha"] = _parse_fecha_mmddyyyy_with_report(df_in[mapping["Fecha"]])
    for k in [k for k in EXPECTED.keys() if k!="Fecha"]:
        out[k] = _to_num(df_in[mapping[k]])
    out = out.dropna(subset=["Fecha"]).sort_values("Fecha").reset_index(drop=True)
    out["mes"] = out["Fecha"].dt.to_period("M").astype(str)
    out["mes_fecha"] = pd.to_datetime(out["Fecha"].dt.to_period("M").astype(str) + "-01", errors="coerce")
    # Producción diaria observada
    out["produccion_t"] = out["rendimiento_real_tph"] * out["tiempo_operativo_real_h/dia"]
    out["produccion_plan_t"] = out["rendimiento_plan_tph"] * out["tiempo_operativo_plan_h/dia"]
    return out

def remove_outliers_iqr(df, cols, factor=1.5):
    d = df.copy()
    mask = pd.Series(True, index=d.index)
    for c in cols:
        q1, q3 = d[c].quantile([0.25, 0.75])
        iqr = q3 - q1
        lo, hi = q1 - factor*iqr, q3 + factor*iqr
        mask &= d[c].between(lo, hi) | d[c].isna()
    return d[mask].copy()

def month_order(values):
    # 'YYYY-MM' -> datetime for sorting
    return sorted(values, key=lambda x: pd.to_datetime(x+"-01"))

def ci95(mean, std, n):
    if n <= 1 or pd.isna(std): return 0.0
    se = std/np.sqrt(n)
    return 1.96*se

def tukey_letters(data, group_col, value_col, alpha=0.05):
    # Compact Letter Display from Tukey HSD
    mc = MultiComparison(data[value_col].values, data[group_col].values)
    res = mc.tukeyhsd(alpha=alpha)
    groups = mc.groupsunique
    # Build adjacency: not significantly different -> connect
    # res.reject True means different; False means same group
    n = len(groups)
    same = {g:set([g]) for g in groups}
    for i in range(len(res.reject)):
        a, b = res._multicomp.pairindices[i]
        ga, gb = groups[a], groups[b]
        if not res.reject[i]:
            same[ga].add(gb); same[gb].add(ga)
    # Greedy letter assignment
    letters = {g:"" for g in groups}
    assigned = []
    for g in groups:
        placed = False
        for lset, letter in assigned:
            # can we put g into this letter group? must be same with all in lset
            if all((gg in same[g]) for gg in lset):
                lset.add(g); letters[g]+=letter; placed=True; break
        if not placed:
            new_letter = chr(ord('A')+len(assigned))
            assigned.append((set([g]), new_letter))
            letters[g]+=new_letter
    return letters  # dict: group -> letters like 'A', 'AB', ...

# ---------- Sidebar (7 controles esenciales) ----------
st.sidebar.title("Controles")

uploaded = st.sidebar.file_uploader("📂 Sube tu CSV", type=["csv","txt"])
sep_semicolon = st.sidebar.checkbox("Separador ';'", value=False)
if uploaded is None:
    st.info("Sube un CSV para iniciar el análisis.")
    st.stop()

try:
    raw = pd.read_csv(uploaded, sep=";" if sep_semicolon else ",")
except Exception:
    raw = pd.read_csv(uploaded, sep=None, engine="python")

df = load_dataset_v2(raw)

# 1) Periodo (mes inicio–fin)
meses_all = month_order(df["mes"].dropna().unique())
mes_ini, mes_fin = st.sidebar.select_slider("Periodo (mes)", options=meses_all, value=(meses_all[0], meses_all[-1]))
mask_periodo = (df["mes"] >= mes_ini) & (df["mes"] <= mes_fin)
dfp = df.loc[mask_periodo].copy()

# 2) Métrica foco
metrica_foco = st.sidebar.radio("Métrica foco", ["Producción","TPH","Horas"], index=0)

# 3) Brecha vs plan (Δ ó %)
modo_brecha = st.sidebar.radio("Brecha vs plan", ["Δ absoluto","Δ %"], index=0)

# 4) Outliers simple (IQR 1.5×)
filtrar_outliers = st.sidebar.checkbox("Quitar outliers (IQR 1.5×)", value=False)
if filtrar_outliers:
    dfp = remove_outliers_iqr(
        dfp,
        ["rendimiento_real_tph","tiempo_operativo_real_h/dia","produccion_t"],
        factor=1.5
    )

# 5) Sensibilidad (rápida)
delta_tph_max = st.sidebar.slider("ΔTPH máximo (%)", 0, 10, 5)
delta_h_max = st.sidebar.slider("ΔHoras máximo (h/d)", 0.0, 4.0, 1.0, step=0.5)
baseline_real = st.sidebar.checkbox("Baseline real (recomendado)", value=True)

# 6) Significancia ANOVA
alpha = 0.05 if st.sidebar.radio("α (significancia ANOVA)", ["0.05","0.10"], index=0)=="0.05" else 0.10

# 7) Exportación (PPT con notas)
incluir_notas = st.sidebar.checkbox("Incluir notas del presentador en PPT", value=True)

st.markdown(f"**Periodo aplicado:** {mes_ini} → {mes_fin}  •  Filas: {len(dfp):,}")

# ================= Resumen KPIs & narrativa =================
def resumen_kpis(d):
    kpis = {
        "Mineral real (t)": d["mineral_procesado_real_t"].sum(),
        "Mineral plan (t)": d["mineral_procesado_plan_t"].sum(),
        "Producción real (t)": d["produccion_t"].sum(),
        "Producción plan (t)": d["produccion_plan_t"].sum(),
        "TPH real prom": d["rendimiento_real_tph"].mean(),
        "TPH plan prom": d["rendimiento_plan_tph"].mean(),
        "Horas reales prom": d["tiempo_operativo_real_h/dia"].mean(),
        "Horas plan prom": d["tiempo_operativo_plan_h/dia"].mean(),
    }
    return pd.DataFrame({"KPI": list(kpis.keys()), "Valor": list(kpis.values())})

kpi_df = resumen_kpis(dfp)
c1, c2 = st.columns([1,1])
with c1:
    st.subheader("KPIs del periodo")
    st.dataframe(kpi_df.style.format({"Valor":"{:,.2f}"}), use_container_width=True)
with c2:
    st.subheader("Narrativa técnica (auto)")
    # Stats base
    tph_mean = dfp["rendimiento_real_tph"].mean()
    tph_med  = dfp["rendimiento_real_tph"].median()
    h_mean   = dfp["tiempo_operativo_real_h/dia"].mean()
    h_med    = dfp["tiempo_operativo_real_h/dia"].median()
    # CV mensual
    cvm = dfp.groupby("mes").agg(
        cv_tph=("rendimiento_real_tph", lambda x: np.std(x, ddof=1)/np.mean(x) if np.mean(x)>0 else np.nan),
        cv_h  =("tiempo_operativo_real_h/dia", lambda x: np.std(x, ddof=1)/np.mean(x) if np.mean(x)>0 else np.nan),
    ).reset_index()
    cv_tph_min, cv_tph_max = np.nanmin(cvm["cv_tph"]), np.nanmax(cvm["cv_tph"])
    cv_h_max = np.nanmax(cvm["cv_h"])
    mes_cv_h_alto = cvm.loc[cvm["cv_h"].idxmax(), "mes"] if cvm["cv_h"].notna().any() else "-"
    # Brechas promedio
    delta_tph_prom = (dfp["rendimiento_real_tph"] - dfp["rendimiento_plan_tph"]).mean()
    delta_h_prom   = (dfp["tiempo_operativo_real_h/dia"] - dfp["tiempo_operativo_plan_h/dia"]).mean()
    # Equivalencia empírica
    sum_TPH = dfp["rendimiento_real_tph"].sum()
    sum_TPH_H = (dfp["rendimiento_real_tph"] * dfp["tiempo_operativo_real_h/dia"]).sum()
    dh_equiv = 0.01 * (sum_TPH_H / max(sum_TPH, 1e-9))

    narrativa = [
        f"Periodo analizado {mes_ini}–{mes_fin}.",
        f"TPH real: media {tph_mean:.1f}, mediana {tph_med:.1f}. Horas reales: media {h_mean:.2f} h/d, mediana {h_med:.2f} h/d.",
        f"Variabilidad mensual: CV(TPH) entre {cv_tph_min:.2f} y {cv_tph_max:.2f}; CV(Horas) máx {cv_h_max:.2f} en {mes_cv_h_alto}.",
        f"Brechas promedio vs plan — ΔTPH {delta_tph_prom:+.1f}, ΔHoras {delta_h_prom:+.1f} h/d.",
        f"Equivalencia empírica: 1% TPH ≈ {dh_equiv:.2f} h/d (si <0.6 favorece disponibilidad; si ≥0.6 favorece capacidad instantánea)."
    ]
    st.markdown("\n\n".join(["- " + s for s in narrativa]))

# ================= Comparación vs plan (mensual) + Brecha mensual =================
st.markdown("---")
st.subheader("Comparación mensual vs plan")

agg_fun = "mean"  # simple y robusto; podrías exponerlo como control
gm = dfp.groupby("mes", as_index=False).agg({
    "rendimiento_real_tph": agg_fun,
    "rendimiento_plan_tph": agg_fun,
    "tiempo_operativo_real_h/dia": agg_fun,
    "tiempo_operativo_plan_h/dia": agg_fun,
    "produccion_t": agg_fun,
    "produccion_plan_t": agg_fun
})
gm = gm.sort_values("mes")

def plot_grouped_bars(dfm, y_real, y_plan, title):
    fig, ax = plt.subplots(figsize=(12,4.8))
    x = np.arange(len(dfm))
    w = 0.38
    ax.bar(x-w/2, dfm[y_real], width=w, label="Real", color=custom_palette[0])
    ax.bar(x+w/2, dfm[y_plan], width=w, label="Plan", color=custom_palette[2])
    ax.set_xticks(x)
    ax.set_xticklabels(dfm["mes"], rotation=45, ha="right")
    ax.set_title(title); ax.legend()
    ax.yaxis.set_major_formatter(StrMethodFormatter("{x:,.0f}"))
    ax.grid(axis="y", alpha=0.25)
    plt.tight_layout()
    return fig

def plot_brecha_mensual(dfm, y_real, y_plan, title, modo="Δ absoluto"):
    d = dfm.copy()
    if modo == "Δ absoluto":
        d["brecha"] = d[y_real] - d[y_plan]
        label_fmt = lambda v: f"{v:+.1f}"
        yfmt = StrMethodFormatter("{x:,.1f}")
    else:
        base = d[y_plan].replace(0, np.nan)
        d["brecha"] = (d[y_real] - d[y_plan]) / base * 100.0
        label_fmt = lambda v: f"{v:+.1f}%"
        yfmt = StrMethodFormatter("{x:,.0f}%")
    fig, ax = plt.subplots(figsize=(12,4.8))
    colors = [custom_palette[0] if v>=0 else "#C64756" for v in d["brecha"]]
    ax.bar(d["mes"], d["brecha"], color=colors)
    ax.set_title(title)
    ax.yaxis.set_major_formatter(yfmt)
    ax.grid(axis="y", alpha=0.25)
    # etiquetas
    for i, v in enumerate(d["brecha"]):
        ax.text(i, v + (0.02*np.nanmax(np.abs(d["brecha"])) if v>=0 else -0.02*np.nanmax(np.abs(d["brecha"]))),
                label_fmt(v), ha="center", va="bottom" if v>=0 else "top", fontsize=9)
    plt.xticks(rotation=45, ha="right"); plt.tight_layout()
    return fig

# Elegir métrica foco para el primer bloque
if metrica_foco == "TPH":
    st.pyplot(plot_grouped_bars(gm, "rendimiento_real_tph", "rendimiento_plan_tph", "TPH — Real vs Plan (mensual)"), use_container_width=True)
    st.pyplot(plot_brecha_mensual(gm, "rendimiento_real_tph", "rendimiento_plan_tph", "Brecha mensual TPH (Real–Plan)", modo_brecha), use_container_width=True)
elif metrica_foco == "Horas":
    st.pyplot(plot_grouped_bars(gm, "tiempo_operativo_real_h/dia", "tiempo_operativo_plan_h/dia", "Horas/día — Real vs Plan (mensual)"), use_container_width=True)
    st.pyplot(plot_brecha_mensual(gm, "tiempo_operativo_real_h/dia", "tiempo_operativo_plan_h/dia", "Brecha mensual Horas (Real–Plan)", modo_brecha), use_container_width=True)
else:
    st.pyplot(plot_grouped_bars(gm, "produccion_t", "produccion_plan_t", "Producción (t) — Real vs Plan (mensual)"), use_container_width=True)
    st.pyplot(plot_brecha_mensual(gm, "produccion_t", "produccion_plan_t", "Brecha mensual Producción (Real–Plan)", modo_brecha), use_container_width=True)

# ================= ANOVA visual (mes) =================
st.markdown("---")
st.subheader("ANOVA visual por mes")

def anova_visual(d, col, ylabel):
    # tabla ANOVA
    model = ols(f"Q('{col}') ~ C(mes)", data=d).fit()
    anova_tbl = sm.stats.anova_lm(model, typ=2).round(4)

    # medias, IC95 y letras Tukey
    stats = d.groupby("mes").agg(mean=(col,"mean"), std=(col,"std"), n=(col,"count")).reset_index()
    stats["ci"] = stats.apply(lambda r: ci95(r["mean"], r["std"], r["n"]), axis=1)

    # Tukey (solo si al menos 2 niveles)
    letters_map = {}
    try:
        letters_map = tukey_letters(d[["mes",col]].dropna(), "mes", col, alpha=alpha)
    except Exception:
        letters_map = {m:"" for m in stats["mes"]}

    stats = stats.sort_values("mes")
    fig, ax = plt.subplots(figsize=(12,4.8))
    x = np.arange(len(stats))
    ax.errorbar(x, stats["mean"], yerr=stats["ci"], fmt="o-", capsize=4, label="Media ± IC95%", color=custom_palette[0])
    ax.set_xticks(x); ax.set_xticklabels(stats["mes"], rotation=45, ha="right")
    ax.set_title(f"{ylabel}: Medias mensuales con IC95% (α={alpha})")
    ax.grid(True, alpha=0.25)
    ax.yaxis.set_major_formatter(StrMethodFormatter("{x:,.1f}"))
    # letras sobre puntos
    for i, row in stats.iterrows():
        ax.text(i, row["mean"] + (0.02*np.nanmax(stats["mean"]) if np.nanmax(stats["mean"])>0 else 0.5),
                letters_map.get(row["mes"], ""), ha="center", va="bottom", fontsize=10, color="#003B5C")
    plt.tight_layout()
    return fig, anova_tbl

# Mostrar dos: TPH y Horas (si foco es Producción igual mostramos ambos para ANOVA)
c1, c2 = st.columns(2)
with c1:
    fig_tph, anova_tph = anova_visual(dfp, "rendimiento_real_tph", "TPH")
    st.pyplot(fig_tph, use_container_width=True)
    st.dataframe(anova_tph, use_container_width=True)
with c2:
    fig_h, anova_h = anova_visual(dfp, "tiempo_operativo_real_h/dia", "Horas/día")
    st.pyplot(fig_h, use_container_width=True)
    st.dataframe(anova_h, use_container_width=True)

# ================= Sensibilidad (iso-%) — baseline real para evitar negativos falsos =================
st.markdown("---")
st.subheader("Sensibilidad (iso-%)")

mean_tph = float(dfp["rendimiento_real_tph"].mean())
mean_h   = float(dfp["tiempo_operativo_real_h/dia"].mean())

if baseline_real:
    # producción media diaria observada
    baseline_iso = float((dfp["produccion_t"].mean()))
else:
    baseline_iso = mean_tph * mean_h

# Construimos malla simple: 0..Δmax con paso 1% y 0.5h
tph_grid = np.arange(0, delta_tph_max+0.001, 1) / 100.0
h_grid = np.arange(0.0, delta_h_max+1e-9, 0.5)

cells = []
for dH in h_grid:
    for dT in tph_grid:
        h_new   = mean_h * (1 + 0) + dH   # sumamos horas absolutas
        tph_new = mean_tph * (1 + dT)
        tons_new = tph_new * h_new
        val = ((tons_new / max(baseline_iso, 1e-9)) - 1.0) * 100.0
        cells.append({"delta_h_abs": dH, "delta_tph_pct": dT, "tons_pct": val})
grid = pd.DataFrame(cells)
pv = grid.pivot(index="delta_tph_pct", columns="delta_h_abs", values="tons_pct")

cmap = mcolors.LinearSegmentedColormap.from_list("afines_heat", ["#E0F7FA", "#00AFAA", "#0B5563"])
fig_iso, ax = plt.subplots(figsize=(10,6))
sns.heatmap(pv, annot=True, fmt=".1f", cmap=cmap, cbar_kws={"label": "Δ Producción (%)"},
            annot_kws={"fontsize": 10}, ax=ax)
ax.set_title("Iso-%: impacto de ΔTPH (%) y ΔHoras (h/d) sobre producción (%)")
ax.set_xlabel("Δ Horas (h/d)"); ax.set_ylabel("Δ TPH (%)")

xt = [f"{v:+.1f}h" for v in pv.columns]
yt = [f"{v*100:.0f}%" for v in pv.index]
# Ticks compat
posx = np.arange(len(xt)) + 0.5
ax.set_xticks(posx); ax.set_xticklabels(xt)
posy = np.arange(len(yt)) + 0.5
ax.set_yticks(posy); ax.set_yticklabels(yt)
ax.tick_params(axis="x", rotation=0); ax.tick_params(axis="y", rotation=0)
plt.tight_layout()
st.pyplot(fig_iso, use_container_width=True)

# ================= PPT con notas del presentador =================
st.markdown("---")
st.subheader("Exportar")

def save_fig_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight")
    buf.seek(0); return buf

def table_fig(df_table, title):
    fig, ax = plt.subplots(figsize=(7.5,3.0))
    ax.axis("off")
    tb = df_table.copy()
    if "sum_sq" in tb.columns:
        tb = tb.round({"sum_sq":4, "df":2, "F":3, "PR(>F)":4})
    tb = tb.reset_index()
    tbl = ax.table(cellText=tb.values, colLabels=tb.columns.tolist(), loc="center")
    tbl.scale(1,1.1)
    ax.set_title(title, pad=6)
    plt.tight_layout()
    return fig

def build_ppt(df_periodo, figs, notes_by_slide, titulo_portada, incluir_notas=True):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_image_slide(title, img_bytes, notes_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(0x32,0x8B,0xA1)
        slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(1.1), width=Inches(12.3))
        if incluir_notas:
            notes = slide.notes_slide.notes_text_frame
            notes.text = notes_text

    # portada
    slide0 = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide0.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = titulo_portada; p.font.size = Pt(36); p.font.bold = True
    p.font.color.rgb = RGBColor(0x32,0x8B,0xA1)
    if incluir_notas:
        slide0.notes_slide.notes_text_frame.text = f"Periodo: {mes_ini}–{mes_fin}. Filas: {len(df_periodo):,}."

    # slides
    for (title, figbytes, notes) in figs:
        add_image_slide(title, figbytes, notes)

    out = io.BytesIO()
    prs.save(out); out.seek(0)
    return out

# Ensamblar figuras clave
figs_export = []

# 1) Comparación mensual y brecha (según métrica foco)
if metrica_foco == "TPH":
    f1 = plot_grouped_bars(gm, "rendimiento_real_tph","rendimiento_plan_tph","TPH — Real vs Plan (mensual)")
    f2 = plot_brecha_mensual(gm, "rendimiento_real_tph","rendimiento_plan_tph","Brecha mensual TPH (Real–Plan)", modo_brecha)
elif metrica_foco == "Horas":
    f1 = plot_grouped_bars(gm, "tiempo_operativo_real_h/dia","tiempo_operativo_plan_h/dia","Horas/día — Real vs Plan (mensual)")
    f2 = plot_brecha_mensual(gm, "tiempo_operativo_real_h/dia","tiempo_operativo_plan_h/dia","Brecha mensual Horas (Real–Plan)", modo_brecha)
else:
    f1 = plot_grouped_bars(gm, "produccion_t","produccion_plan_t","Producción (t) — Real vs Plan (mensual)")
    f2 = plot_brecha_mensual(gm, "produccion_t","produccion_plan_t","Brecha mensual Producción (Real–Plan)", modo_brecha)

figs_export.append(("Comparación mensual vs plan", save_fig_bytes(f1),
                    "Comparación por mes, medias. Revisar meses con sobre/subcumplimiento para ajustar plan."))
figs_export.append(("Brecha mensual vs plan", save_fig_bytes(f2),
                    "Brechas en Δ ó %. Identificar top/bottom 3 meses y causas operativas."))

# 2) ANOVA visual + tablas
figs_export.append(("ANOVA visual — TPH", save_fig_bytes(fig_tph),
                    "Medias mensuales TPH con IC95%. Letras Tukey indican grupos no diferentes."))
figs_export.append(("ANOVA visual — Horas", save_fig_bytes(fig_h),
                    "Medias mensuales Horas con IC95%. Interpretar meses ‘altos/bajos’."))
figs_export.append(("ANOVA — tabla TPH", save_fig_bytes(table_fig(anova_tph, "ANOVA — TPH")),
                    "Si p<α hay efecto mes; considerar replanificación en meses críticos."))
figs_export.append(("ANOVA — tabla Horas", save_fig_bytes(table_fig(anova_h, "ANOVA — Horas")),
                    "Tamaño de efecto: usar η² (opcional) para priorizar."))

# 3) Sensibilidad iso-%
figs_export.append(("Sensibilidad iso-%", save_fig_bytes(fig_iso),
                    "Mapa con baseline real. Aumentos en TPH/horas deben mostrar Δ producción positivo."))

titulo = f"Chancado — Consolidado ({mes_ini} a {mes_fin})"
ppt_bytes = build_ppt(dfp, figs_export, None, titulo, incluir_notas=incluir_notas)

st.download_button(
    "⬇️ Descargar PPTX 16:9",
    data=ppt_bytes,
    file_name=f"chancado_fusion_{pd.Timestamp.now():%Y%m%d_%H%M}.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

